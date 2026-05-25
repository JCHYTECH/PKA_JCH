from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_DECK = Path(
    "JCH_Inbox/03_PROJECTS/04_NUANCES/docs/nuances_b2b_mission_deck_art_direction_v5_premium.pptx"
)

SLIDE_W = 13.333
SLIDE_H = 7.5

BG = RGBColor(248, 247, 243)
NAVY = RGBColor(26, 43, 60)
TEAL = RGBColor(24, 115, 115)
SAND = RGBColor(227, 220, 207)
AMBER = RGBColor(190, 128, 56)
DARK = RGBColor(44, 44, 44)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(98, 107, 117)
BLACK_TERRAIN = RGBColor(13, 11, 9)
BRUN_NUIT = RGBColor(42, 36, 32)

FONT_SANS = "Space Grotesk"
FONT_SERIF = "Lora"
FONT_MONO = "IBM Plex Mono"


SLIDES = [
    {
        "type": "cover",
        "title": "Structurer la vente B2B\nsans diluer la promesse Nuances",
        "subtitle": "Mission B2B | Nuances",
        "meta": "Mai 2026 | Version de travail",
    },
    {
        "type": "content",
        "section": "01",
        "title": "Définir Nuances",
        "statement": "Nuances n'est pas un simple installateur.\nC'est un opérateur d'optimisation énergétique B2B.",
        "label": "Leviers clés",
        "bullets": [
            "Données, EMS, stockage et financement dans une même logique de valeur.",
            "Pas un catalogue PV + batterie : une architecture de performance énergétique.",
            "Objectif : transformer le site client en actif piloté, mesurable et rentable.",
        ],
        "footer": "Pas un installateur — un copilote énergie",
    },
    {
        "type": "content",
        "section": "02",
        "title": "Vision B2B",
        "variant": "aurel",
        "statement": "Construire une activité B2B qui vend une performance énergétique mesurable, pas seulement des équipements.",
        "label": "Ambition",
        "bullets": [
            "Réduire les coûts énergétiques et les émissions de CO2.",
            "Créer des revenus additionnels via pilotage, stockage et flexibilité.",
            "Devenir le copilote énergie du client avec une relation récurrente et défendable.",
        ],
        "footer": "Donnée + technique + financement + preuve économique",
    },
    {
        "type": "content",
        "section": "03",
        "title": "Rôle du duo Arteon",
        "variant": "aurel",
        "statement": "Arteon intervient comme binôme de transformation commerciale : clarifier, structurer, accélérer.",
        "label": "Notre rôle",
        "bullets": [
            "Donner une lecture stratégique claire du marché, de l'offre et des arbitrages.",
            "Transformer cette lecture en processus technico-commercial concret, outillé et pilotable.",
            "Accompagner Nuances sur le terrain jusqu'à validation réelle des hypothèses de vente.",
        ],
        "footer": "Vision stratégique + exécution technico-commerciale",
    },
    {
        "type": "content",
        "section": "04",
        "title": "Contexte marché BESS",
        "statement": "Un marché en forte croissance, plus structuré, plus lisible, avec des segments encore mal couverts.",
        "label": "Marché et évolution",
        "bullets": [
            "Europe : 18,1 Md$ en 2024 -> 87,34 Md$ projetés en 2033 ; +45 % de capacité installée en 2025.",
            "LFP s'impose comme standard du stationnaire ; le stockage devient un outil économique de pilotage.",
            "Plus de 8 000 communautés d'énergie dans l'UE : un cadre qui favorise l'autoconsommation et les ressources distribuées.",
        ],
        "footer": "Sources : BESS-acteurs | BESS-technologies | BESS-legal",
    },
    {
        "type": "content",
        "section": "05",
        "title": "Pourquoi structurer maintenant",
        "statement": "Nuances a déjà les actifs. Il lui manque la machine commerciale capable de les monétiser proprement en B2B.",
        "label": "Frottements actuels",
        "bullets": [
            "Prospection trop large et rendez-vous insuffisamment préparés.",
            "Commercial, technique et financement restent insuffisamment alignés.",
            "Le B2B exige une preuve de rentabilité rapide et un cycle mieux piloté que le B2C.",
        ],
        "footer": "Structurer avant d'accélérer",
    },
    {
        "type": "content",
        "section": "06",
        "title": "Besoin central de la mission",
        "statement": "La mission doit faire gagner du temps en amont, augmenter la qualité des rendez-vous et accélérer la signature.",
        "label": "Transformation attendue",
        "bullets": [
            "Mieux qualifier les prospects et réduire les rendez-vous inutiles.",
            "Préparer des hypothèses de solution et de rentabilité avant visite.",
            "Créer un enchaînement vente -> technique -> offre -> exécution piloté par KPI.",
        ],
        "footer": "Passer plus vite du lead au contrat",
    },
    {
        "type": "content",
        "section": "07",
        "title": "Cibles prioritaires",
        "statement": "PME industrielles familiales avec consommation significative et chaîne de décision courte.",
        "label": "Pourquoi cette cible",
        "bullets": [
            "Potentiel PV / BESS / EMS réel et douleur économique visible.",
            "Cabine MT ou HT : filtre simple pour exclure rapidement les sites peu pertinents.",
            "Cycles plus courts et accès décisionnel plus simple que dans les grands comptes structurés.",
        ],
        "footer": "Commencer étroit pour signer plus vite",
    },
    {
        "type": "content",
        "section": "08",
        "title": "Offres B2B à articuler",
        "statement": "Ne pas vendre un catalogue. Vendre une progression de valeur : diagnostiquer, mesurer, piloter, équiper, financer, optimiser.",
        "label": "Briques Nuances",
        "bullets": [
            "Audit énergie, boîtier / capteurs / passerelle, Nuances Brain / EMS.",
            "BESS industriel, PV / Nuances Sun, contrats énergie et valorisation.",
            "Financement : leasing, tiers investisseur et logique de cash-flow.",
        ],
        "footer": "Une chaîne de valeur, pas une juxtaposition d'offres",
    },
    {
        "type": "grid4",
        "section": "09",
        "title": "Processus technico-commercial cible",
        "items": [
            ("01-02", "Cibler\n& présélectionner", "Filtre MT/HT\nScore prospect\nGo / No Go"),
            ("03-04", "Qualifier\n& préparer", "Données minimales\nDossier vendeur\nHypothèse solution"),
            ("05-07", "Rendez-vous\nData\nScénario", "Visite terrain\nCollecte de données\nModélisation"),
            ("08-10", "Proposer\nSigner\nSuivre", "Offre chiffrée\nNégociation\nMonitoring"),
        ],
        "footer": "Le vendeur arrive préparé, pas avec une brochure",
    },
    {
        "type": "grid4",
        "section": "10",
        "title": "6 leviers pour raccourcir la signature",
        "items": [
            ("01", "Mieux\ncibler", "Go / No Go\ndès le 1er contact"),
            ("02", "Exiger les\ndonnées clés", "Contrat fourniture\nPuissance cabine\nConsommation"),
            ("03", "Préparer le\ndossier vendeur", "Hypothèse solution\nEstimation gains"),
            ("04-06", "Financement\nROI\nScénarios", "Achat / leasing tôt\nVendre sur cash-flow\nMinimal / optimal / financé"),
        ],
        "footer": "Moins de rendez-vous inutiles, plus de crédibilité",
    },
    {
        "type": "content",
        "section": "11",
        "title": "Dossier préliminaire avant rendez-vous",
        "statement": "Le commercial arrive avec une hypothèse de solution, une estimation de rentabilité et une liste de validations.",
        "label": "Contenu attendu",
        "bullets": [
            "Identité prospect, contexte énergétique, score d'opportunité.",
            "Hypothèses connues, opportunités par offre et ordre de grandeur des gains.",
            "Scénarios minimal / optimal / financé, questions à poser et prochaines validations techniques.",
        ],
        "footer": "Une base solide avant le premier rendez-vous",
    },
    {
        "type": "content",
        "section": "12",
        "title": "Boîte à outils vente et pilotage",
        "statement": "Sans outils, la stratégie reste théorique. Avec outils, elle devient exécutable et mesurable.",
        "label": "Outils clés",
        "bullets": [
            "Vente : grille de qualification, questionnaire, dossier vendeur, argumentaires.",
            "Technique : modèle d'audit, calculateur ROI / cash-flow / CO2, guide de collecte.",
            "Pilotage : CRM, pipeline, dossiers générés, suivi installation et performance.",
        ],
        "footer": "Des outils simples à lire et simples à mettre à jour",
    },
    {
        "type": "content",
        "section": "13",
        "title": "KPI de pilotage",
        "statement": "Prospection / Conversion / Performance — les KPI qui pilotent réellement la machine B2B.",
        "label": "Familles KPI",
        "bullets": [
            "Prospection : taux de leads qualifiés, temps par lead, taux de rendez-vous obtenus.",
            "Conversion : taux rendez-vous -> audit -> offre -> signature.",
            "Performance : ROI moyen, économies estimées, revenus BESS / flexibilité, CO2 évité.",
        ],
        "footer": "Piloter la vitesse, la qualité et la rentabilité",
    },
    {
        "type": "grid4",
        "section": "14",
        "title": "Phases de mission proposées",
        "items": [
            ("01", "Cadrage\nstratégique", "Clarifier offre B2B\nCible prioritaire\nPoints ouverts"),
            ("02", "Qualification\n& présélection", "Grille scoring\nGo / No Go\nQuestionnaire"),
            ("03", "Dossier\nvendeur", "Dossier préliminaire\nHypothèses\nScénarios solution"),
            ("04-06", "Modèle éco\nCRM\nDéploiement", "Financement\nKPI & CRM\nFormation terrain"),
        ],
        "footer": "La mission produit un système de vente B2B",
    },
    {
        "type": "content",
        "section": "15",
        "title": "Périmètre de l'engagement",
        "statement": "4 à 5 jours par mois par consultant — feuille de temps, frais et validation terrain progressive.",
        "label": "Cadre de mission",
        "bullets": [
            "Présence sur site, protocole technico-commercial et itérations à partir des visites prospects.",
            "Nuances fournit accès, documents, formations produits et interlocuteur opérationnel.",
            "Périmètre élargi possible : lobbying ciblé, projets de contrats et négociations associées.",
        ],
        "footer": "Mission outillante, pas mission décorative",
    },
    {
        "type": "content",
        "section": "16",
        "title": "Conditions de succès",
        "statement": "Le risque principal n'est pas l'absence d'idée. C'est de produire de la documentation sans transformer l'exécution.",
        "label": "Conditions clés",
        "bullets": [
            "Alignement rapide sur l'offre prioritaire et arbitrages clairs sur la cible.",
            "Accès aux données minimales et implication terrain d'un spécialiste Nuances.",
            "Intégration du processus dans le CRM et pilotage régulier par KPI.",
        ],
        "footer": "La mission doit changer la vitesse de qualification",
    },
    {
        "type": "content",
        "section": "17",
        "title": "Ce que Nuances résout déjà",
        "statement": "Nuances répond déjà à plusieurs blocages du marché BESS là où les grands acteurs sont peu agiles.",
        "label": "Risques déjà mitigés",
        "bullets": [
            "CAPEX élevé et ROI flou : audit, boîtier, preuve économique et financement intégré.",
            "Prospection non qualifiée : filtre cabine MT/HT, ciblage PME et logique Go / No Go.",
            "Offre trop complexe : Nuances vend une performance pilotée, pas une batterie isolée.",
        ],
        "footer": "Données + pilotage + financement + service",
    },
    {
        "type": "content",
        "section": "18",
        "title": "Points à vérifier avant de démarrer",
        "statement": "Ces points ne bloquent pas la mission. Ils conditionnent la réussite réelle de l'activité B2B.",
        "label": "Risques à cadrer",
        "bullets": [
            "Rôle exact de Nuances, conformité, certification et passeport batterie 2027.",
            "Transformation B2C vers B2B, structure du leasing et service post-installation.",
            "Partenaires OEM, assurabilité, sécurité incendie et choix du segment initial.",
        ],
        "footer": "Choisir le bon rôle avant de scaler",
    },
    {
        "type": "grid4",
        "section": "19",
        "title": "Recommandations immédiates",
        "items": [
            ("01", "Choisir le\nsegment", "Démarrer sur un périmètre\nétroit et rentable"),
            ("02", "Fixer le rôle\nde Nuances", "Revendeur, intégrateur,\nopérateur ou financeur"),
            ("03", "Sécuriser les\npartenaires", "Techno certifiée\nFinance\nAssurance"),
            ("04", "Lancer des\npilotes", "3 à 5 pilotes\nstrictement sélectionnés\navant extension"),
        ],
        "footer": "Choisir -> cadrer -> tester -> scaler",
    },
    {
        "type": "closing",
        "title": "Nuances peut gagner sa place en B2B\nsi elle choisit le bon segment,\nle bon rôle et le bon niveau de service.",
        "subtitle": "Le marché existe. La technologie est mûre.\nLa réussite dépend maintenant de la focalisation et de l'exécution.",
    },
]


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=20, color=DARK, bold=False,
                align=PP_ALIGN.LEFT, font_name=FONT_SANS, valign=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, fill_color, line_color=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    return shape


def add_content_slide(prs, data):
    if data.get("variant") == "aurel":
        return add_aurel_slide(prs, data)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_rect(slide, 0, 0, 0.26, SLIDE_H, TEAL)
    add_rect(slide, 0.42, 0.55, 12.1, 0.05, SAND)
    add_textbox(slide, 0.6, 0.25, 1.0, 0.4, data["section"], 11, color=AMBER, font_name=FONT_MONO)
    add_textbox(slide, 0.9, 0.9, 6.2, 0.8, data["title"], 24, color=NAVY, bold=True, font_name=FONT_SANS)
    add_textbox(slide, 0.9, 1.7, 5.7, 1.35, data["statement"], 18, color=DARK, font_name=FONT_SERIF)

    add_rect(slide, 7.05, 1.0, 5.4, 4.5, WHITE, line_color=SAND, radius=True)
    add_textbox(slide, 7.35, 1.25, 4.8, 0.45, data["label"], 13, color=AMBER, bold=True, font_name=FONT_SANS)

    y = 1.85
    for bullet in data["bullets"]:
        add_rect(slide, 7.38, y + 0.07, 0.1, 0.1, TEAL, radius=True)
        add_textbox(slide, 7.62, y, 4.45, 0.7, bullet, 13.5, color=DARK, font_name=FONT_SANS)
        y += 1.0

    add_rect(slide, 0.9, 6.35, 11.6, 0.65, NAVY)
    add_textbox(slide, 1.15, 6.49, 11.0, 0.3, data["footer"], 11.5, color=WHITE, font_name=FONT_MONO)
    return slide


def add_grid4_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_rect(slide, 0, 0, 0.26, SLIDE_H, TEAL)
    add_rect(slide, 0.42, 0.55, 12.1, 0.05, SAND)
    add_textbox(slide, 0.6, 0.25, 1.0, 0.4, data["section"], 11, color=AMBER, font_name=FONT_MONO)
    add_textbox(slide, 0.9, 0.9, 7.0, 0.7, data["title"], 24, color=NAVY, bold=True, font_name=FONT_SANS)

    cards = [
        (0.9, 1.8),
        (3.95, 1.8),
        (7.0, 1.8),
        (10.05, 1.8),
    ]
    for (left, top), item in zip(cards, data["items"]):
        num, title, desc = item
        add_rect(slide, left, top, 2.65, 3.2, WHITE, line_color=SAND, radius=True)
        add_textbox(slide, left + 0.2, top + 0.2, 0.6, 0.35, num, 11, color=AMBER, bold=True, font_name=FONT_MONO)
        add_textbox(slide, left + 0.2, top + 0.65, 2.1, 0.8, title, 17, color=NAVY, bold=True, font_name=FONT_SANS)
        add_textbox(slide, left + 0.2, top + 1.55, 2.05, 1.25, desc, 12.5, color=DARK, font_name=FONT_SANS)

    add_rect(slide, 0.9, 6.35, 11.6, 0.65, NAVY)
    add_textbox(slide, 1.15, 6.49, 11.0, 0.3, data["footer"], 11.5, color=WHITE, font_name=FONT_MONO)
    return slide


def add_cover_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK_TERRAIN)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BLACK_TERRAIN)
    add_rect(slide, 0, 0, 0.18, SLIDE_H, AMBER)
    add_rect(slide, 9.9, 0, 3.43, SLIDE_H, BRUN_NUIT)
    add_textbox(slide, 10.05, 0.48, 2.15, 0.3, "Arteon", 20, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name=FONT_SERIF)
    add_textbox(slide, 0.75, 0.55, 2.2, 0.25, "ARTEON MASTER", 10, color=WHITE, bold=True, font_name=FONT_SANS)
    add_textbox(slide, 0.8, 1.45, 6.0, 1.9, data["title"], 27, color=WHITE, bold=True, font_name=FONT_SANS)
    add_rect(slide, 0.82, 4.2, 1.65, 0.03, AMBER)
    add_textbox(slide, 0.82, 4.45, 3.8, 0.4, "Jean-Claude Havaux x Ghislain Lacrosse", 11.5, color=AMBER, font_name=FONT_SANS)
    add_textbox(slide, 0.8, 6.0, 2.6, 0.18, "Nuances B2B", 10, color=WHITE, font_name=FONT_SANS)
    add_textbox(slide, 0.8, 6.34, 3.8, 0.3, "Mission B2B | " + data["subtitle"].replace("Mission B2B | ", ""), 10.5, color=MUTED, font_name=FONT_MONO)
    add_textbox(slide, 10.25, 5.95, 2.55, 0.5, data["meta"], 10.5, color=WHITE, font_name=FONT_MONO)
    return slide


def add_closing_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_rect(slide, 0, 6.55, SLIDE_W, 0.95, BLACK_TERRAIN)
    add_textbox(slide, 0.8, 0.55, 2.2, 0.25, "ARTEON MASTER", 10, color=DARK, bold=True, font_name=FONT_SANS)
    add_textbox(slide, 10.95, 0.42, 1.3, 0.32, "Arteon", 20, color=BLACK_TERRAIN, align=PP_ALIGN.RIGHT, font_name=FONT_SERIF)
    add_rect(slide, 0.78, 1.0, 1.55, 0.03, AMBER)
    add_textbox(slide, 0.8, 1.4, 6.1, 1.95, data["title"], 25, color=BLACK_TERRAIN, bold=True, font_name=FONT_SANS)
    add_textbox(slide, 0.82, 4.18, 5.7, 0.85, data["subtitle"], 17.5, color=MUTED, font_name=FONT_SERIF)
    add_rect(slide, 9.35, 1.2, 2.45, 2.45, SAND)
    add_rect(slide, 9.83, 1.68, 1.5, 1.5, AMBER)
    add_textbox(slide, 0.9, 6.88, 3.8, 0.2, "Nuances B2B | Synthèse finale", 10.5, color=WHITE, font_name=FONT_MONO)
    return slide


def add_aurel_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_textbox(slide, 0.95, 0.8, 1.2, 0.22, f"Section {data['section']}", 10.5, color=AMBER, font_name=FONT_SANS)
    add_textbox(slide, 10.95, 0.68, 1.3, 0.3, "Arteon", 20, color=BLACK_TERRAIN, align=PP_ALIGN.RIGHT, font_name=FONT_SERIF)
    add_rect(slide, 0.95, 1.15, 1.7, 0.03, AMBER)
    add_textbox(slide, 0.95, 1.55, 5.6, 1.25, data["title"], 27, color=BLACK_TERRAIN, bold=True, font_name=FONT_SERIF, italic=True)
    add_textbox(slide, 0.98, 3.0, 5.5, 1.1, data["statement"], 18, color=DARK, font_name=FONT_SERIF)
    add_textbox(slide, 0.98, 4.55, 5.1, 0.28, data["footer"], 11.5, color=AMBER, font_name=FONT_MONO)
    add_rect(slide, 8.45, 0.95, 3.35, 3.35, SAND)
    add_rect(slide, 8.93, 1.43, 2.38, 2.38, BRUN_NUIT)
    add_rect(slide, 9.38, 1.88, 1.48, 1.48, AMBER)
    add_textbox(slide, 8.65, 4.72, 3.3, 0.25, data["label"], 10.5, color=MUTED, font_name=FONT_SANS)
    y = 5.0
    for bullet in data["bullets"]:
        add_textbox(slide, 8.65, y, 3.7, 0.42, bullet, 11.5, color=DARK, font_name=FONT_SANS)
        y += 0.48
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    for data in SLIDES:
        if data["type"] == "cover":
            add_cover_slide(prs, data)
        elif data["type"] == "closing":
            add_closing_slide(prs, data)
        elif data["type"] == "grid4":
            add_grid4_slide(prs, data)
        else:
            add_content_slide(prs, data)

    prs.save(str(OUTPUT_DECK))


if __name__ == "__main__":
    build()
